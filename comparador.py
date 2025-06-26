import streamlit as st
import pandas as pd
from pptx import Presentation
from io import BytesIO
import openpyxl

def extract_chart_data_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    charts = []
    for i, slide in enumerate(prs.slides):
        slide_title = None
        for shape in slide.shapes:
            if shape.has_text_frame and not slide_title:
                text = shape.text.strip()
                if text.lower().startswith("indicador"):
                    slide_title = text
            if shape.has_chart:
                chart = shape.chart
                data = []
                categories = [c.label for c in chart.plots[0].categories]
                for series in chart.series:
                    name = series.name
                    values = series.values
                    row = [name] + list(values)
                    data.append(row)
                df = pd.DataFrame(data, columns=["Identificador"] + categories)
                charts.append((slide_title or f"Diapositiva {i+1}", df))
    return charts


def extract_blocks_by_slide_marker(excel_path, marker_prefix="## Diapositiva"):
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    blocks = []

    for sheet in wb.sheetnames:
        ws = wb[sheet]
        current_marker = None
        current_data = []
        for row in ws.iter_rows(values_only=True):
            if row and isinstance(row[0], str) and row[0].strip().lower().startswith(marker_prefix.lower()):
                if current_marker and current_data:
                    df = pd.DataFrame(current_data[1:], columns=current_data[0])
                    blocks.append((current_marker, df))
                current_marker = row[0].strip()
                current_data = []
            elif current_marker:
                if all(cell is None for cell in row):
                    if current_data:
                        df = pd.DataFrame(current_data[1:], columns=current_data[0])
                        blocks.append((current_marker, df))
                        current_marker = None
                        current_data = []
                else:
                    current_data.append(list(row))
        if current_marker and current_data:
            df = pd.DataFrame(current_data[1:], columns=current_data[0])
            blocks.append((current_marker, df))
    return blocks


def normalize_dataframe(df):
    df = df.copy()
    if df.columns[0].lower() in ["marca", "identificador"]:
        df = df.set_index(df.columns[0])
    else:
        df = df.T
        df.columns = df.iloc[0]
        df = df[1:]
        df.index.name = "Identificador"
    df = df.sort_index().sort_index(axis=1)
    return df

def compare_dataframes_flexibly(df1, df2):
    df1_norm = normalize_dataframe(df1)
    df2_norm = normalize_dataframe(df2)

    differences = []
    for row_label in df1_norm.index:
        for col_label in df1_norm.columns:
            val1 = df1_norm.at[row_label, col_label] if row_label in df1_norm.index and col_label in df1_norm.columns else None
            val2 = df2_norm.at[row_label, col_label] if row_label in df2_norm.index and col_label in df2_norm.columns else None
            if pd.isna(val1) and pd.isna(val2):
                continue
            if val1 != val2:
                differences.append({
                    "Identificador": row_label,
                    "Categoría": col_label,
                    "Valor PPT": val1,
                    "Valor Excel": val2
                })
    return differences

def main():
    st.title("Comparador de Gráficos PowerPoint vs Excel (con marcadores)")

    pptx_file = st.file_uploader("Carga tu archivo PowerPoint (.pptx)", type="pptx")
    excel_file = st.file_uploader("Carga tu archivo Excel (.xlsx)", type="xlsx")

    if pptx_file and excel_file:
        pptx_charts = extract_chart_data_from_pptx(pptx_file)
        excel_blocks = extract_blocks_from_excel_by_marker(excel_file, markers=["indicador 1", "indicador 2"])

        all_differences = []

        for slide_title, chart_df in pptx_charts:
            match_found = False
            for block_title, excel_df in excel_blocks:
                if slide_title.lower() == block_title.lower():
                    differences = compare_dataframes_flexibly(chart_df, excel_df)
                    if not differences:
                        st.success(f"✅ {slide_title} coincide con el bloque '{block_title}' del Excel.")
                    else:
                        st.error(f"❌ {slide_title} tiene diferencias con el bloque '{block_title}':")
                        st.dataframe(pd.DataFrame(differences))
                        for diff in differences:
                            diff["Gráfico"] = slide_title
                            diff["Bloque Excel"] = block_title
                        all_differences.extend(differences)
                    match_found = True
                    break
            if not match_found:
                st.warning(f"⚠️ No se encontró un bloque en Excel que coincida con el marcador '{slide_title}'.")
            with st.expander(f"Ver datos del gráfico: {slide_title}"):
                st.dataframe(chart_df)

        if all_differences:
            st.markdown("### 📥 Descargar diferencias como Excel")
            output = BytesIO()
            pd.DataFrame(all_differences).to_excel(output, index=False, engine='openpyxl')
            st.download_button(
                label="Descargar diferencias_ppt_vs_excel.xlsx",
                data=output.getvalue(),
                file_name="diferencias_ppt_vs_excel.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )

if __name__ == "__main__":
    main()
